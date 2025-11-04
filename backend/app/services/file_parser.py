import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import io
from typing import Tuple, List


class FileParser:
    """Parse various file formats and extract text"""
    
    @staticmethod
    async def parse_pdf(file_content: bytes) -> Tuple[str, List[str]]:
        """Extract text and topics from PDF"""
        try:
            doc = fitz.open(stream=file_content, filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text()
            
            # Extract potential topics (simple heuristic: capitalized phrases)
            topics = FileParser._extract_topics(text)
            return text, topics
        except Exception as e:
            raise Exception(f"Error parsing PDF: {str(e)}")
    
    @staticmethod
    async def parse_docx(file_content: bytes) -> Tuple[str, List[str]]:
        """Extract text and topics from DOCX"""
        try:
            doc = Document(io.BytesIO(file_content))
            text = "\n".join([para.text for para in doc.paragraphs])
            topics = FileParser._extract_topics(text)
            return text, topics
        except Exception as e:
            raise Exception(f"Error parsing DOCX: {str(e)}")
    
    @staticmethod
    async def parse_pptx(file_content: bytes) -> Tuple[str, List[str]]:
        """Extract text and topics from PPTX"""
        try:
            prs = Presentation(io.BytesIO(file_content))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            topics = FileParser._extract_topics(text)
            return text, topics
        except Exception as e:
            raise Exception(f"Error parsing PPTX: {str(e)}")
    
    @staticmethod
    async def parse_image(file_content: bytes) -> Tuple[str, List[str]]:
        """Extract text from image using OCR"""
        try:
            image = Image.open(io.BytesIO(file_content))
            text = pytesseract.image_to_string(image)
            topics = FileParser._extract_topics(text)
            return text, topics
        except Exception as e:
            # If OCR fails, return empty text
            return "", []
    
    @staticmethod
    def _extract_topics(text: str) -> List[str]:
        """Extract potential topics from text (optimized for speed)"""
        # Limit text processing for large documents
        if len(text) > 50000:
            text = text[:50000]  # Only process first 50k characters
        
        topics = []
        lines = text.split('\n')
        
        # Process only first 200 lines for speed
        for line in lines[:200]:
            line = line.strip()
            # Look for lines that might be headings (short, capitalized)
            if 5 < len(line) < 80 and line[0].isupper():
                # Skip lines with too many special characters
                if sum(c.isalnum() or c.isspace() for c in line) / len(line) > 0.7:
                    topics.append(line)
                    if len(topics) >= 20:  # Early exit
                        break
        
        # Return unique topics, limited to 20
        return list(dict.fromkeys(topics))[:20]  # Faster than set for preserving order
